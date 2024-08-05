import os
import pdfplumber
from docx import Document
from pptx import Presentation

"""

CONVERT ALL INPUT FILES TO MD FORMAT

"""
def convert_pdf_to_md(pdf_path, output_md_path):
    with pdfplumber.open(pdf_path) as pdf:
        with open(output_md_path, 'w', encoding='utf-8') as md_file:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    md_file.write(text + '\n\n')

def convert_docx_to_md(docx_path, output_md_path):
    doc = Document(docx_path)
    with open(output_md_path, 'w', encoding='utf-8') as md_file:
        for para in doc.paragraphs:
            md_file.write(para.text + '\n\n')

def convert_pptx_to_md(pptx_path, output_md_path):
    prs = Presentation(pptx_path)
    with open(output_md_path, 'w', encoding='utf-8') as md_file:
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    md_file.write(paragraph.text + '\n\n')

def convert_file_to_md(file_path, output_folder):
    file_ext = os.path.splitext(file_path)[1].lower()
    file_name = os.path.splitext(os.path.basename(file_path))[0]
    os.makedirs(output_folder, exist_ok=True)
    output_md_path = os.path.join(output_folder, file_name + '.md')

    if file_ext == '.pdf':
        convert_pdf_to_md(file_path, output_md_path)
    elif file_ext == '.docx':
        convert_docx_to_md(file_path, output_md_path)
    elif file_ext == '.pptx':
        convert_pptx_to_md(file_path, output_md_path)
    else:
        print(f"Unsupported file format: {file_ext}")
    print(f"Output saved to {output_md_path}")

def process_folder(input_directory, output_directory):
    supported_extensions = ['.pdf', '.docx', '.pptx']
    for filename in os.listdir(input_directory):
        if os.path.splitext(filename)[1].lower() in supported_extensions and not filename.startswith('~$'):
            file_path = os.path.join(input_directory, filename)
            convert_file_to_md(file_path, output_directory)
        else:
            print(f"Ignored temporary or unsupported file: {filename}")
def main():
    INPUT_PATH = 'raw_data/test_data/ECON1220'
    OUTPUT_PATH = 'md_data'  
    process_folder(INPUT_PATH, OUTPUT_PATH)

if __name__ == "__main__":
    main()